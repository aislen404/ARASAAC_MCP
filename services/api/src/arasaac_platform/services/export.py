import base64
import html
import json
import zipfile
from collections.abc import Awaitable, Callable
from io import BytesIO

import httpx
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen.canvas import Canvas

from arasaac_platform.domain.materials import Material, MaterialStatus
from arasaac_platform.governance.license import validate_material_license
from arasaac_platform.schemas.materials import ExportManifest


class ExportBlockedError(ValueError):
    """Raised when review or compliance does not permit export."""


ImageFetcher = Callable[[str], Awaitable[bytes]]


def export_html(material: Material) -> bytes:
    _ensure_exportable(material)
    blocks = "\n".join(_render_html_block(block) for block in material.blocks)
    document = f"""<!doctype html>
<html lang="es">
<head><meta charset="utf-8"><title>{html.escape(material.title)}</title></head>
<body>
<main>
<h1>{html.escape(material.title)}</h1>
<ol>{blocks}</ol>
<footer>{html.escape(material.attribution_text)}</footer>
</main>
</body>
</html>"""
    return document.encode()


def export_docx(material: Material) -> bytes:
    _ensure_exportable(material)
    document = Document()
    document.add_heading(material.title, level=1)
    for block in material.blocks:
        if block.pictogram is None:
            continue
        document.add_paragraph(block.text)
        document.add_paragraph(f"Pictograma ARASAAC: {block.pictogram.label}")
    document.add_paragraph(material.attribution_text)
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def export_pptx(material: Material) -> bytes:
    _ensure_exportable(material)
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = material.title
    for block in material.blocks:
        if block.pictogram is None:
            continue
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = block.pictogram.label
        body = slide.placeholders[1]
        body.text = block.text
        body.text_frame.paragraphs[0].font.size = Pt(18)
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def export_zip_package(material: Material) -> bytes:
    _ensure_exportable(material)
    manifest = ExportManifest(
        material_id=str(material.material_id),
        version=material.version,
        attribution=material.attribution_text,
        pictograms=material.pictograms,
        human_review_approved=True,
    )
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("material.html", export_html(material))
        archive.writestr(
            "manifest.json",
            json.dumps(manifest.model_dump(mode="json"), ensure_ascii=False, indent=2),
        )
        archive.writestr("ATTRIBUTION.txt", material.attribution_text)
    return buffer.getvalue()


async def export_pdf(
    material: Material,
    fetch_image: ImageFetcher | None = None,
) -> bytes:
    _ensure_exportable(material)
    fetcher = fetch_image or _fetch_official_image
    buffer = BytesIO()
    canvas = Canvas(buffer, pagesize=A4)
    width, height = A4
    y = height - 54
    canvas.setTitle(material.title)
    canvas.setFont("Helvetica-Bold", 18)
    canvas.drawString(54, y, material.title)
    y -= 48

    for block in material.blocks:
        if block.pictogram is None:
            continue
        image = await fetcher(str(block.pictogram.source_url))
        canvas.drawImage(
            ImageReader(BytesIO(image)),
            54,
            y - 72,
            width=72,
            height=72,
            preserveAspectRatio=True,
            mask="auto",
        )
        canvas.setFont("Helvetica", 12)
        canvas.drawString(144, y - 36, block.text)
        y -= 96
        if y < 140:
            canvas.showPage()
            y = height - 54

    canvas.setFont("Helvetica", 8)
    text = canvas.beginText(54, 72)
    for line in _wrap(material.attribution_text, 100):
        text.textLine(line)
    canvas.drawText(text)
    canvas.save()
    return buffer.getvalue()


def encode_export(content: bytes) -> str:
    return base64.b64encode(content).decode("ascii")


def _render_html_block(block: object) -> str:
    from arasaac_platform.domain.materials import MaterialBlock

    if not isinstance(block, MaterialBlock) or block.pictogram is None:
        return ""
    return (
        '<li class="material-item">'
        f'<img src="{html.escape(str(block.pictogram.source_url))}" '
        f'alt="{html.escape(block.pictogram.label)}">'
        f"<p>{html.escape(block.text)}</p>"
        "</li>"
    )


def _ensure_exportable(material: Material) -> None:
    if material.status != MaterialStatus.APPROVED:
        raise ExportBlockedError("Solo se puede exportar material aprobado.")
    compliance = validate_material_license(material)
    if not compliance.valid:
        raise ExportBlockedError("; ".join(compliance.errors))


async def _fetch_official_image(url: str) -> bytes:
    if not url.startswith("https://static.arasaac.org/pictograms/"):
        raise ExportBlockedError("Origen de imagen no permitido.")
    async with httpx.AsyncClient(follow_redirects=False, timeout=8.0) as client:
        response = await client.get(url)
        response.raise_for_status()
        return response.content


def _wrap(text: str, width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if len(candidate) > width and current:
            lines.append(current)
            current = word
        else:
            current = candidate
    if current:
        lines.append(current)
    return lines
